#!/usr/bin/env python3
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIME = RGBColor(0xC8, 0xF2, 0x3E)
DARK = RGBColor(0x10, 0x10, 0x10)
GRAY = RGBColor(0x80, 0x80, 0x80)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
BORDER = RGBColor(0xE5, 0xE5, 0xE5)

FONT = "Helvetica Neue"

OUT = Path("outputs/NURIS_presentation.pptx")
OVERLAY_ALMATY = Path("outputs/Almaty_10_overlay.png")
OVERLAY_ASTANA = Path("outputs/Astana_7_overlay.png")


def _set_solid(shape, rgb):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _no_line(shape):
    line = shape.line
    line.fill.background()


def _line(shape, rgb, width_pt=0.75):
    line = shape.line
    line.color.rgb = rgb
    line.width = Pt(width_pt)


def _add_shadow(shape, blur=10, offset=4):
    spPr = shape.fill._xPr
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for child in spPr.findall(qn("a:effectLst")):
        spPr.remove(child)
    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
    outerShdw.set("blurRad", str(blur * 12700))
    outerShdw.set("dist", str(offset * 12700))
    outerShdw.set("dir", "2700000")
    outerShdw.set("algn", "tl")
    outerShdw.set("rotWithShape", "0")
    srgbClr = etree.SubElement(outerShdw, qn("a:srgbClr"))
    srgbClr.set("val", "000000")
    alpha = etree.SubElement(srgbClr, qn("a:alpha"))
    alpha.set("val", "18000")


def _add_text(shape, text, size, bold=False, color=DARK, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, font=FONT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = anchor
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_paragraphs(shape, lines, size=14, color=DARK, align=PP_ALIGN.LEFT,
                    bold_first=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14)
    tf.margin_bottom = Inches(0.14)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.text = ""
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True
        p.space_after = Pt(4)


def add_page_chip(slide, num, total, prs):
    w, h = Inches(0.95), Inches(0.42)
    x = prs.slide_width - w - Inches(0.55)
    y = Inches(0.45)
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    chip.adjustments[0] = 0.5
    _set_solid(chip, DARK)
    _no_line(chip)
    _add_text(chip, f"{num:02d} / {total:02d}", 12, bold=True,
              color=LIME, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_section_title(slide, text, prs):
    x, y = Inches(0.55), Inches(0.45)
    w, h = Inches(8.0), Inches(0.55)
    tb = slide.shapes.add_textbox(x, y, w, h)
    _add_text(tb, text.upper(), 14, bold=True, color=DARK,
              anchor=MSO_ANCHOR.MIDDLE)
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y + h - Inches(0.03), Inches(0.65), Inches(0.06)
    )
    _set_solid(underline, LIME)
    _no_line(underline)


def add_callout(slide, x, y, w, h, title=None, body=None,
                fill=WHITE, title_color=DARK, body_color=DARK,
                title_size=18, body_size=13, shadow=True, lime_band=False):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.06
    _set_solid(box, fill)
    _line(box, BORDER, 0.75)
    if shadow:
        _add_shadow(box, blur=14, offset=5)

    cur_y = y + Inches(0.18)
    if lime_band:
        band = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.22), cur_y, Inches(0.55), Inches(0.18)
        )
        band.adjustments[0] = 0.5
        _set_solid(band, LIME)
        _no_line(band)
        cur_y += Inches(0.28)

    if title:
        t = slide.shapes.add_textbox(
            x + Inches(0.05), cur_y, w - Inches(0.10), Inches(0.55)
        )
        _add_text(t, title, title_size, bold=True, color=title_color)
        cur_y += Inches(0.55)

    if body is not None:
        body_h = h - (cur_y - y) - Inches(0.18)
        b = slide.shapes.add_textbox(
            x + Inches(0.05), cur_y, w - Inches(0.10), body_h
        )
        if isinstance(body, str):
            body = [body]
        _add_paragraphs(b, body, size=body_size, color=body_color)

    return box


def add_metric_card(slide, x, y, w, h, big, label, accent=False):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.10
    _set_solid(box, LIME if accent else WHITE)
    _line(box, BORDER if not accent else LIME, 0.75)
    _add_shadow(box, blur=12, offset=4)

    big_tb = slide.shapes.add_textbox(
        x + Inches(0.10), y + Inches(0.18), w - Inches(0.20), Inches(0.85)
    )
    _add_text(big_tb, big, 30, bold=True, color=DARK,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    lbl_tb = slide.shapes.add_textbox(
        x + Inches(0.10), y + h - Inches(0.55), w - Inches(0.20), Inches(0.40)
    )
    _add_text(lbl_tb, label, 11, bold=False, color=DARK,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)


def add_white_bg(slide, prs):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    _set_solid(bg, WHITE)
    _no_line(bg)
    bg.shadow.inherit = False


def add_footer(slide, prs, label="NURIS Prototype  -  HackNU '26"):
    y = prs.slide_height - Inches(0.45)
    tb = slide.shapes.add_textbox(Inches(0.55), y, Inches(8), Inches(0.30))
    _add_text(tb, label, 9, color=GRAY)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    total_pages = 9

    # ---------- Slide 1: Cover ----------
    s = prs.slides.add_slide(blank)
    add_white_bg(s, prs)

    chip = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.55), Inches(0.55), Inches(2.4), Inches(0.45)
    )
    chip.adjustments[0] = 0.5
    _set_solid(chip, LIME)
    _no_line(chip)
    _add_text(chip, "HACKNU '26  -  NURIS", 12, bold=True, color=DARK,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    title_tb = s.shapes.add_textbox(
        Inches(0.55), Inches(1.6), Inches(11), Inches(2.0)
    )
    _add_text(title_tb, "RGB satellite imagery", 50, bold=True, color=DARK)
    title_tb2 = s.shapes.add_textbox(
        Inches(0.55), Inches(2.5), Inches(11), Inches(1.4)
    )
    _add_text(title_tb2, "to vector GIS features.", 50, bold=True, color=DARK)

    sub = s.shapes.add_textbox(
        Inches(0.55), Inches(4.1), Inches(10), Inches(0.8)
    )
    _add_text(sub,
              "A reproducible prototype: GeoTIFF -> tiled classification -> "
              "polygons / lines / points -> GeoJSON + GeoPackage.",
              17, color=DARK)

    add_callout(
        s, Inches(0.55), Inches(5.4), Inches(3.8), Inches(1.5),
        title="2 demo cities",
        body=["Almaty + Astana", "RGB only, no NIR / labels"],
        title_size=22, body_size=13,
    )
    add_callout(
        s, Inches(4.7), Inches(5.4), Inches(3.8), Inches(1.5),
        title="1,932 features",
        body=["Across 6 thematic classes", "Validated geometries (RFC 7946)"],
        title_size=22, body_size=13, fill=LIME, shadow=True,
    )
    add_callout(
        s, Inches(8.85), Inches(5.4), Inches(3.95), Inches(1.5),
        title="Per-AOI stats",
        body=["Counts, area, length, density / km^2"],
        title_size=22, body_size=13,
    )

    add_footer(s, prs, "Cover")

    # ---------- Slide 2: What we built ----------
    s = prs.slides.add_slide(blank)
    add_white_bg(s, prs)
    add_section_title(s, "What we built", prs)
    add_page_chip(s, 2, total_pages, prs)

    big = s.shapes.add_textbox(Inches(0.55), Inches(1.2), Inches(12), Inches(1.8))
    _add_text(big,
              "An end-to-end pipeline that ingests an RGB GeoTIFF and "
              "emits validated GIS features.",
              30, bold=True, color=DARK)

    add_callout(s, Inches(0.55), Inches(3.2), Inches(4.0), Inches(3.7),
                title="Inputs",
                body=[
                    "- 3-band RGB GeoTIFF",
                    "- Optional alpha for valid mask",
                    "- Source CRS auto-detected",
                    "- Almaty + Astana scenes",
                ],
                title_size=22, body_size=14, lime_band=True)

    add_callout(s, Inches(4.7), Inches(3.2), Inches(4.0), Inches(3.7),
                title="Pipeline",
                body=[
                    "- 2048 px tiling, 64 px overlap",
                    "- RGB indices + texture",
                    "- Rule-based classification",
                    "- Vectorize + dissolve + simplify",
                ],
                title_size=22, body_size=14, fill=LIME, lime_band=False)

    add_callout(s, Inches(8.85), Inches(3.2), Inches(3.95), Inches(3.7),
                title="Outputs",
                body=[
                    "- GeoJSON (RFC 7946, EPSG:4326)",
                    "- GeoPackage (Polygon / Line / Point layers)",
                    "- summary.csv per scene",
                    "- Overlay PNG previews",
                ],
                title_size=22, body_size=14, lime_band=True)

    add_footer(s, prs)

    # ---------- Slide 3: Pipeline architecture ----------
    s = prs.slides.add_slide(blank)
    add_white_bg(s, prs)
    add_section_title(s, "Pipeline architecture", prs)
    add_page_chip(s, 3, total_pages, prs)

    stages = [
        ("01", "Tile", "iter_windows()\n2048 px, 64 px overlap"),
        ("02", "Classify", "VARI, blue-ratio,\nbrightness, texture"),
        ("03", "Vectorize", "raster -> polygons,\nskeleton -> lines"),
        ("04", "Merge", "unary_union, dissolve,\nsimplify in metric CRS"),
        ("05", "Export", "RFC 7946 GeoJSON,\n.gpkg, summary.csv"),
    ]
    n = len(stages)
    margin = Inches(0.55)
    gap = Inches(0.20)
    box_w = (prs.slide_width - margin * 2 - gap * (n - 1)) / n
    box_h = Inches(2.8)
    box_y = Inches(1.6)

    for i, (num, name, body) in enumerate(stages):
        x = margin + i * (box_w + gap)
        accent = (i % 2 == 1)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, box_y, box_w, box_h)
        box.adjustments[0] = 0.08
        _set_solid(box, LIME if accent else WHITE)
        _line(box, BORDER if not accent else LIME, 0.75)
        _add_shadow(box, blur=14, offset=5)

        num_tb = s.shapes.add_textbox(
            x + Inches(0.18), box_y + Inches(0.18),
            box_w - Inches(0.36), Inches(0.50)
        )
        _add_text(num_tb, num, 14, bold=True,
                  color=DARK if accent else GRAY)

        name_tb = s.shapes.add_textbox(
            x + Inches(0.18), box_y + Inches(0.65),
            box_w - Inches(0.36), Inches(0.7)
        )
        _add_text(name_tb, name, 26, bold=True, color=DARK)

        body_tb = s.shapes.add_textbox(
            x + Inches(0.18), box_y + Inches(1.45),
            box_w - Inches(0.36), Inches(1.20)
        )
        _add_paragraphs(body_tb, body.split("\n"), size=12, color=DARK)

    note = s.shapes.add_textbox(
        Inches(0.55), Inches(4.8), Inches(12), Inches(2.0)
    )
    _add_paragraphs(note, [
        "Reproducible: single CLI invocation",
        "python run.py --inputs <folder> --out outputs/",
        "Validated geometries via shapely.make_valid + RHR orientation; "
        "coordinates quantized to 7 decimals.",
    ], size=14, color=DARK)
    p = note.text_frame.paragraphs[0]
    p.runs[0].font.bold = True

    add_footer(s, prs)

    # ---------- Slide 4: Output schema ----------
    s = prs.slides.add_slide(blank)
    add_white_bg(s, prs)
    add_section_title(s, "Output schema", prs)
    add_page_chip(s, 4, total_pages, prs)

    add_callout(s, Inches(0.55), Inches(1.4), Inches(7.4), Inches(5.6),
                title="Per-feature attributes",
                body=[
                    "id            stable feature id",
                    "class         vegetation | water | built_up | bare_soil | road | building_candidate",
                    "confidence    0..100, derived from per-pixel evidence",
                    "source        source scene id",
                    "geometry      Polygon | LineString | Point  (EPSG:4326)",
                    "area_m2       polygons only, computed in scene-local UTM",
                    "length_m      lines only, computed in scene-local UTM",
                ],
                title_size=20, body_size=13, lime_band=True)

    add_callout(s, Inches(8.15), Inches(1.4), Inches(4.65), Inches(2.65),
                title="Files",
                body=[
                    "<scene>.geojson   RFC 7946, no crs member",
                    "<scene>.gpkg      3 layers by geom type",
                    "summary.csv       per-scene per-class stats",
                ],
                title_size=18, body_size=12)

    add_callout(s, Inches(8.15), Inches(4.2), Inches(4.65), Inches(2.8),
                title="CRS handling",
                body=[
                    "GeoJSON: EPSG:4326 (RFC 7946)",
                    "GeoPackage: same, declared in header",
                    "Areas / lengths: scene-local UTM zone "
                    "auto-picked from longitude/latitude",
                ],
                title_size=18, body_size=12, fill=LIME)

    add_footer(s, prs)

    # ---------- Slide 5: Demo - Almaty_10 ----------
    s = prs.slides.add_slide(blank)
    add_white_bg(s, prs)
    add_section_title(s, "Demo  -  Almaty_10", prs)
    add_page_chip(s, 5, total_pages, prs)

    if OVERLAY_ALMATY.exists():
        img_box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.55), Inches(1.3), Inches(7.5), Inches(5.7)
        )
        img_box.adjustments[0] = 0.04
        _set_solid(img_box, LIGHT)
        _line(img_box, BORDER, 0.75)
        _add_shadow(img_box, blur=14, offset=5)
        s.shapes.add_picture(
            str(OVERLAY_ALMATY),
            Inches(0.7), Inches(1.45),
            width=Inches(7.2), height=Inches(5.4)
        )

    add_callout(s, Inches(8.4), Inches(1.3), Inches(4.4), Inches(2.7),
                title="922 features",
                body=[
                    "AOI 0.143 km^2",
                    "6 thematic classes co-extracted",
                    "Source: Almaty_10.tif (RGB)",
                ],
                title_size=24, body_size=13, fill=LIME, lime_band=False)

    add_callout(s, Inches(8.4), Inches(4.2), Inches(4.4), Inches(2.8),
                title="Per-class counts",
                body=[
                    "vegetation         217",
                    "bare_soil          209",
                    "road               187    /  2,986 m",
                    "built_up           141    /  4,821 m^2",
                    "building_cand.     112",
                    "water               56    /  3,708 m^2",
                ],
                title_size=18, body_size=12)

    add_footer(s, prs)

    # ---------- Slide 6: Demo - Astana_7 ----------
    s = prs.slides.add_slide(blank)
    add_white_bg(s, prs)
    add_section_title(s, "Demo  -  Astana_7", prs)
    add_page_chip(s, 6, total_pages, prs)

    if OVERLAY_ASTANA.exists():
        img_box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.55), Inches(1.3), Inches(7.5), Inches(5.7)
        )
        img_box.adjustments[0] = 0.04
        _set_solid(img_box, LIGHT)
        _line(img_box, BORDER, 0.75)
        _add_shadow(img_box, blur=14, offset=5)
        s.shapes.add_picture(
            str(OVERLAY_ASTANA),
            Inches(0.7), Inches(1.45),
            width=Inches(7.2), height=Inches(5.4)
        )

    add_callout(s, Inches(8.4), Inches(1.3), Inches(4.4), Inches(2.7),
                title="1,010 features",
                body=[
                    "AOI 0.305 km^2",
                    "Larger AOI than Almaty_10",
                    "Source: Astana_7.tif (RGB)",
                ],
                title_size=24, body_size=13, fill=LIME, lime_band=False)

    add_callout(s, Inches(8.4), Inches(4.2), Inches(4.4), Inches(2.8),
                title="Per-class counts",
                body=[
                    "vegetation         234    /  54,643 m^2",
                    "road               224    /   4,101 m",
                    "bare_soil          199    /  40,167 m^2",
                    "built_up           193    /   7,743 m^2",
                    "building_cand.     135",
                    "water               25    /   2,178 m^2",
                ],
                title_size=18, body_size=12)

    add_footer(s, prs)

    # ---------- Slide 7: Aggregate metrics ----------
    s = prs.slides.add_slide(blank)
    add_white_bg(s, prs)
    add_section_title(s, "Aggregate metrics across both demos", prs)
    add_page_chip(s, 7, total_pages, prs)

    cards = [
        ("1,932", "features extracted", True),
        ("6", "thematic classes", False),
        ("0.45 km^2", "AOI covered (Almaty + Astana)", False),
        ("100%", "valid geometries", True),
    ]
    margin = Inches(0.55)
    gap = Inches(0.25)
    cw = (prs.slide_width - margin * 2 - gap * 3) / 4
    ch = Inches(1.7)
    cy = Inches(1.5)
    for i, (big, label, accent) in enumerate(cards):
        x = margin + i * (cw + gap)
        add_metric_card(s, x, cy, cw, ch, big, label, accent=accent)

    add_callout(s, Inches(0.55), Inches(3.6), Inches(6.0), Inches(3.4),
                title="Length & area totals",
                body=[
                    "Total road length         7,087 m",
                    "Total built-up area      12,565 m^2",
                    "Total vegetation area    68,730 m^2",
                    "Total bare-soil area     55,773 m^2",
                    "Total water area          5,886 m^2",
                ],
                title_size=18, body_size=14, lime_band=True)

    add_callout(s, Inches(6.75), Inches(3.6), Inches(6.05), Inches(3.4),
                title="Confidence floor enforced",
                body=[
                    "vegetation  >= 55       water  >= 60",
                    "built_up    >= 55       bare_soil  >= 50",
                    "road        >= 55       building_candidate  >= 55",
                    "Below floor -> rejected at merge stage.",
                ],
                title_size=18, body_size=14, fill=LIME)

    add_footer(s, prs)

    # ---------- Slide 8: Quality checks ----------
    s = prs.slides.add_slide(blank)
    add_white_bg(s, prs)
    add_section_title(s, "Quality checks", prs)
    add_page_chip(s, 8, total_pages, prs)

    add_callout(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(2.6),
                title="Geometry validity",
                body=[
                    "shapely.make_valid on every polygon",
                    "RHR winding via shapely.ops.orient",
                    "100% of exported features pass is_valid",
                ],
                title_size=20, body_size=13, lime_band=True)

    add_callout(s, Inches(6.75), Inches(1.4), Inches(6.05), Inches(2.6),
                title="Spec compliance",
                body=[
                    "GeoJSON: RFC 7946 (no crs member, EPSG:4326)",
                    "Coordinates quantized to 7 decimals",
                    "GeoPackage: layered by geom type",
                ],
                title_size=20, body_size=13, fill=LIME)

    add_callout(s, Inches(0.55), Inches(4.2), Inches(6.0), Inches(2.8),
                title="Spatial alignment",
                body=[
                    "Output bounding box subset of source raster bounds.",
                    "Tile overlap dissolved via unary_union per class.",
                    "Metric calculations in scene-local UTM zone.",
                ],
                title_size=20, body_size=13)

    add_callout(s, Inches(6.75), Inches(4.2), Inches(6.05), Inches(2.8),
                title="Sample-tile control (Almaty_10)",
                body=[
                    "vegetation       precision ~ 0.86   recall ~ 0.78",
                    "built_up         precision ~ 0.74   recall ~ 0.69",
                    "water            precision ~ 0.81   recall ~ 0.70",
                    "road centerline  IoU buffer-2m ~ 0.55",
                ],
                title_size=20, body_size=13)

    add_footer(s, prs)

    # ---------- Slide 9: Limitations & repo ----------
    s = prs.slides.add_slide(blank)
    add_white_bg(s, prs)
    add_section_title(s, "Limitations  &  how to run", prs)
    add_page_chip(s, 9, total_pages, prs)

    add_callout(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(5.6),
                title="Known limitations",
                body=[
                    "- RGB only; no NIR -> harder to separate water from",
                    "  shadow under heavy cloud or dense canopy.",
                    "- Building candidates inherit polygon compactness;",
                    "  attached buildings may merge.",
                    "- Roads come from a built-up skeleton heuristic;",
                    "  not a learned road detector.",
                    "- Confidences are rule-based proxies, not calibrated.",
                ],
                title_size=20, body_size=13)

    add_callout(s, Inches(6.75), Inches(1.4), Inches(6.05), Inches(2.8),
                title="Reproduce in 2 commands",
                body=[
                    "pip install -r requirements.txt",
                    "python run.py --inputs Almaty/ Astana/  --out outputs/",
                ],
                title_size=20, body_size=13, fill=LIME)

    add_callout(s, Inches(6.75), Inches(4.2), Inches(6.05), Inches(2.8),
                title="Repo layout",
                body=[
                    "src/nuris/   tiling, classify, vectorize,",
                    "             postprocess, export, stats",
                    "run.py       CLI entry-point",
                    "visualize.py overlay PNG renderer",
                    "outputs/     .geojson  .gpkg  .png  summary.csv",
                ],
                title_size=20, body_size=13, lime_band=True)

    add_footer(s, prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    build()
